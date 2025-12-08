import json
import logging
import os
from io import BytesIO
from pathlib import Path
from functools import lru_cache
from typing import Any, Callable, Dict, List, Optional, Tuple

import copy
import fitz
import pytesseract
import requests
from PIL import Image, ImageDraw, ImageFont
from docx import Document
from openpyxl import load_workbook
from openpyxl.comments import Comment
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app import config_manager
from app.ollama_utils import generate_with_metrics

LOGGER = logging.getLogger("insightaudio.translator")
SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".xlsx", ".pdf"}


@lru_cache(maxsize=128)
def get_model_tuning(model: str) -> Dict[str, Any]:
    cfg = config_manager.get_config()
    tuning = (cfg.get("MODEL_TUNING") or {}).get(model, {})
    return copy.deepcopy(tuning) if tuning else {}


def merge_model_options(model: str, defaults: Dict[str, Any]) -> Dict[str, Any]:
    options = dict(defaults)
    tuning = get_model_tuning(model)
    tuning_options = dict(tuning.get("options") or {})
    context_length = tuning.get("context_length")
    if context_length and "num_ctx" not in tuning_options:
        tuning_options["num_ctx"] = context_length
    options.update(tuning_options)
    return options


class TranslationBatch:
    """Помогает собирать и применять пакетный перевод для разных форматов."""

    def __init__(self, scope: str):
        self.scope = scope
        self._blocks: List[Dict[str, Any]] = []

    def add(self, text: str, apply_fn: Callable[[str], None]) -> None:
        text = (text or "").strip()
        if not text:
            return
        block_id = f"{self.scope}-{len(self._blocks)}"
        self._blocks.append({"id": block_id, "text": text, "apply": apply_fn})

    @property
    def blocks(self) -> List[Dict[str, Any]]:
        return self._blocks

    def payload(self) -> List[Dict[str, str]]:
        return [{"id": b["id"], "text": b["text"]} for b in self._blocks]

    def apply(self, translations: Dict[str, str], fallback: Callable[[str], str]) -> None:
        for block in self._blocks:
            translated = translations.get(block["id"])
            if not translated and fallback:
                translated = fallback(block["text"])
            if translated is not None:
                block["apply"](translated)


def translate_document(
    input_path: str,
    target_language: str,
    model: Optional[str] = None,
    backend: Optional[str] = None,
    custom_api_url: Optional[str] = None,
    pdf_reflow: bool = False,
    image_mode: str = "notes",
    translation_mode: str = "block",
) -> str:
    """
    Универсальная точка входа для перевода документов.
    В зависимости от расширения запускает профиль для DOCX/PPTX/XLSX/PDF.
    Каждый профиль разбивает документ на минимальные смысловые блоки
    и синхронно вызывает LLM для каждого блока, сохраняя исходную структуру.
    """
    cfg = config_manager.get_config()
    model_name = model or cfg.get("DEFAULT_TRANSLATE_MODEL", cfg.get("DEFAULT_SUMMARIZE_MODEL"))
    backend_name = backend or cfg.get("TRANSLATE_BACKEND", "ollama")
    max_ratio = float(cfg.get("TRANSLATION_MAX_RATIO", 1.35))
    ext = Path(input_path).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Формат {ext} не поддерживается. Доступны: {', '.join(sorted(SUPPORTED_EXTENSIONS))}")

    output_path = _build_output_path(input_path)
    chunk_mode = translation_mode if translation_mode in {"block", "full"} else cfg.get("DEFAULT_TRANSLATION_MODE", "block")
    if ext == ".pdf" and chunk_mode == "full":
        LOGGER.info("Режим единого запроса пока недоступен для PDF, используется блочное преобразование.")
        chunk_mode = "block"

    translator = _DocTranslator(
        backend=backend_name,
        model=model_name,
        target_lang=target_language,
        max_ratio=max_ratio,
        custom_api_url=custom_api_url,
        image_mode=image_mode or cfg.get("DEFAULT_IMAGE_TRANSLATION_MODE", "notes"),
        chunk_mode=chunk_mode,
    )

    if ext == ".docx":
        translator.translate_docx(input_path, output_path)
    elif ext == ".pptx":
        translator.translate_pptx(input_path, output_path)
    elif ext == ".xlsx":
        translator.translate_xlsx(input_path, output_path)
    elif ext == ".pdf":
        translator.translate_pdf(input_path, output_path, reflow=pdf_reflow)
    return output_path


def _build_output_path(input_path: str) -> str:
    directory, filename = os.path.split(input_path)
    stem, ext = os.path.splitext(filename)
    return os.path.join(directory, f"{stem}_translated{ext}")


class _DocTranslator:
    def __init__(
        self,
        backend: str,
        model: str,
        target_lang: str,
        max_ratio: float,
        custom_api_url: Optional[str] = None,
        image_mode: str = "notes",
        chunk_mode: str = "block",
    ):
        cfg = config_manager.get_config()
        self.backend = backend
        self.model = model
        self.target_lang = target_lang
        self.max_ratio = max_ratio
        self.custom_api_url = custom_api_url
        self.image_mode = image_mode or "notes"
        self.chunk_mode = chunk_mode if chunk_mode in {"block", "full"} else "block"
        self.image_notes: List[str] = []
        self.ocr_lang = cfg.get("OCR_LANG", "rus+eng")
        if backend == "ollama":
            self.base_url = (cfg.get("OLLAMA_API_BASE") or "http://localhost:11434").rstrip("/")
        elif backend == "llama_cpp":
            self.base_url = cfg.get("LLAMA_CPP_API_URL") or "http://localhost:8080"
        elif backend == "custom_api":
            self.base_url = custom_api_url or cfg.get("CUSTOM_SUMMARY_API_URL") or ""
        else:
            self.base_url = (cfg.get("OLLAMA_API_BASE") or "http://localhost:11434").rstrip("/")
        """
        Вспомогательный объект, который инкапсулирует пайплайны перевода каждого формата.
        chunk_mode позволяет переключаться между блочным и монолитным способом общения с LLM.
        """

    def translate_docx(self, input_path: str, output_path: str) -> None:
        """Перевод DOCX целиком: обход параграфов, таблиц, inline-изображений, сбор заметок."""
        LOGGER.info("Перевод DOCX: %s", input_path)
        self.image_notes = []
        doc = Document(input_path)
        if self.chunk_mode == "full":
            batch = TranslationBatch("docx")
            for paragraph in doc.paragraphs:
                self._collect_docx_paragraph(batch, paragraph)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            self._collect_docx_paragraph(batch, paragraph)
            translations = self._translate_batch(batch, "docx")
            batch.apply(translations, self._fallback_translate_single)
        else:
            for paragraph in doc.paragraphs:
                self._translate_paragraph(paragraph)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            self._translate_paragraph(paragraph)
        self._process_docx_images(doc)
        self._append_image_notes_docx(doc)
        doc.save(output_path)

    def _collect_docx_paragraph(self, batch: TranslationBatch, paragraph) -> None:
        text = (paragraph.text or "").strip()
        if not text:
            return
        batch.add(text, lambda value, target=paragraph: self._apply_docx_paragraph(target, value))

    def _apply_docx_paragraph(self, paragraph, text: str) -> None:
        paragraph.text = text

    def _collect_ppt_shape_text(self, batch: TranslationBatch, shape) -> None:
        if shape.has_text_frame:
            original = shape.text
            if original and original.strip():
                batch.add(original, lambda value, frame=shape.text_frame: self._replace_text_preserving_format(frame, value))
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                self._collect_ppt_shape_text(batch, sub_shape)

    def _translate_batch(self, batch: TranslationBatch, scope: str) -> Dict[str, str]:
        if not batch.blocks:
            return {}
        instruction = (
            "Тебе передан JSON-массив объектов вида {\"id\": \"...\", \"text\": \"...\"}. "
            f"Переведи поле text на {self.target_lang}. Верни JSON-массив с теми же id "
            "и полем \"translation\" (строка). Не добавляй комментариев и пояснений."
        )
        payload = json.dumps(batch.payload(), ensure_ascii=False)
        try:
            raw = self._call_translate_api(instruction, payload)
        except Exception as exc:
            LOGGER.warning("Не удалось выполнить пакетный перевод (%s): %s", scope, exc)
            return {}
        try:
            data = json.loads(raw)
        except json.JSONDecodeError:
            LOGGER.warning("Пакетный перевод (%s): нераспознанный ответ: %s", scope, raw[:200])
            return {}
        mapping: Dict[str, str] = {}
        if isinstance(data, list):
            for item in data:
                if not isinstance(item, dict):
                    continue
                block_id = item.get("id")
                translated = item.get("translation") or item.get("text")
                if block_id and isinstance(translated, str):
                    mapping[block_id] = translated.strip()
        return mapping

    def _fallback_translate_single(self, text: str) -> str:
        translated, _ = self._translate_text(text)
        return translated

    def translate_pptx(self, input_path: str, output_path: str) -> None:
        """Перевод PPTX: текстовые фреймы, группы, заметки, изображения на каждом слайде."""
        LOGGER.info("Перевод PPTX: %s", input_path)
        prs = Presentation(input_path)
        for idx, slide in enumerate(prs.slides, start=1):
            self.image_notes = []
            notes_text_parts = []
            if self.chunk_mode == "full":
                batch = TranslationBatch(f"pptx-{idx}")
                for shape in slide.shapes:
                    self._collect_ppt_shape_text(batch, shape)
                translations = self._translate_batch(batch, f"pptx-{idx}")
                batch.apply(translations, self._fallback_translate_single)
            else:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        original = shape.text
                        if original.strip():
                            translated, comment = self._translate_text(original)
                            self._replace_text_preserving_format(shape.text_frame, translated)
                            if comment:
                                notes_text_parts.append(comment)
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        for sub_shape in shape.shapes:
                            if sub_shape.has_text_frame:
                                original = sub_shape.text
                                if original.strip():
                                    translated, comment = self._translate_text(original)
                                    self._replace_text_preserving_format(sub_shape.text_frame, translated)
                                    if comment:
                                        notes_text_parts.append(comment)
            self._process_ppt_pictures(slide.shapes, idx, notes_text_parts)
            
            # Добавляем комментарии в заметки слайда
            if notes_text_parts:
                self._add_notes_to_slide(slide, "\n\n".join(notes_text_parts))
        
        prs.save(output_path)
    
    def _process_ppt_pictures(self, shapes, slide_idx: int, notes_text_parts: List[str]) -> None:
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blob = shape.image.blob
                new_blob, note = self._process_image_bytes(blob, f"Слайд {slide_idx}")
                if new_blob is not None and self.image_mode == "redesign":
                    shape.image.replace(new_blob)
                if note:
                    notes_text_parts.append(note)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._process_ppt_pictures(shape.shapes, slide_idx, notes_text_parts)

    def _replace_text_preserving_format(self, text_frame, new_text: str) -> None:
        """Заменяет текст в text_frame, сохраняя форматирование каждого run и paragraph"""
        if not text_frame.paragraphs:
            return
        
        # Сохраняем форматирование первого paragraph (выравнивание, отступы и т.д.)
        first_para = text_frame.paragraphs[0]
        para_format = {
            'alignment': first_para.alignment,
            'level': first_para.level,
            'space_after': first_para.space_after,
            'space_before': first_para.space_before,
        }
        
        # Собираем все форматирования из существующих runs
        format_styles = []
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                # Сохраняем размер шрифта в Points, если он задан
                font_size = None
                if run.font.size:
                    try:
                        # Сохраняем как есть (уже в Points)
                        font_size = run.font.size
                    except (AttributeError, TypeError):
                        pass
                
                style = {
                    'font_name': run.font.name,
                    'font_size': font_size,
                    'bold': run.font.bold,
                    'italic': run.font.italic,
                    'underline': run.font.underline,
                    'color': run.font.color,
                }
                format_styles.append(style)
        
        # Если нет runs, создаем стиль по умолчанию из первого paragraph
        if not format_styles and text_frame.paragraphs:
            if first_para.runs:
                first_run = first_para.runs[0]
                font_size = None
                if first_run.font.size:
                    try:
                        font_size = first_run.font.size
                    except (AttributeError, TypeError):
                        pass
                default_style = {
                    'font_name': first_run.font.name,
                    'font_size': font_size,
                    'bold': first_run.font.bold,
                    'italic': first_run.font.italic,
                    'underline': first_run.font.underline,
                    'color': first_run.font.color,
                }
            else:
                # Используем стиль paragraph (если доступен)
                default_style = {}
                if hasattr(first_para, 'font'):
                    font_size = None
                    if first_para.font.size:
                        try:
                            font_size = first_para.font.size
                        except (AttributeError, TypeError):
                            pass
                    default_style = {
                        'font_name': first_para.font.name,
                        'font_size': font_size,
                        'bold': first_para.font.bold,
                        'italic': first_para.font.italic,
                        'underline': first_para.font.underline,
                        'color': first_para.font.color,
                    }
            if default_style:
                format_styles = [default_style]
        
        # Очищаем все paragraphs
        text_frame.clear()
        
        # Добавляем новый paragraph с переведенным текстом
        paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
        
        # Восстанавливаем форматирование paragraph
        if para_format.get('alignment') is not None:
            paragraph.alignment = para_format['alignment']
        if para_format.get('level') is not None:
            paragraph.level = para_format['level']
        if para_format.get('space_after') is not None:
            paragraph.space_after = para_format['space_after']
        if para_format.get('space_before') is not None:
            paragraph.space_before = para_format['space_before']
        
        # Используем первое форматирование для всего текста
        if format_styles:
            style = format_styles[0]
            run = paragraph.add_run()
            run.text = new_text
            
            # Применяем сохраненное форматирование шрифта
            if style.get('font_name'):
                try:
                    run.font.name = style['font_name']
                except (AttributeError, TypeError):
                    pass
            if style.get('font_size'):
                try:
                    # Применяем размер шрифта (уже в Points)
                    run.font.size = style['font_size']
                except (AttributeError, TypeError):
                    pass
            if style.get('bold') is not None:
                try:
                    run.font.bold = style['bold']
                except (AttributeError, TypeError):
                    pass
            if style.get('italic') is not None:
                try:
                    run.font.italic = style['italic']
                except (AttributeError, TypeError):
                    pass
            if style.get('underline') is not None:
                try:
                    run.font.underline = style['underline']
                except (AttributeError, TypeError):
                    pass
            if style.get('color'):
                try:
                    run.font.color = style['color']
                except (AttributeError, TypeError):
                    pass
        else:
            # Если нет форматирования, просто добавляем текст
            paragraph.text = new_text
    
    def _add_notes_to_slide(self, slide, notes_text: str) -> None:
        """Добавляет текст в заметки слайда (notes)"""
        try:
            # Получаем notes_slide (всегда существует в python-pptx)
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            
            # Проверяем, что notes_text_frame существует
            if notes_text_frame is None:
                LOGGER.warning("notes_text_frame не доступен для слайда")
                return
            
            # Добавляем текст в заметки
            existing_text = notes_text_frame.text or ""
            if existing_text.strip():
                notes_text_frame.text = existing_text + "\n\n" + notes_text
            else:
                notes_text_frame.text = notes_text
        except Exception as e:
            LOGGER.warning("Не удалось добавить заметки к слайду: %s", e)

    def _process_image_bytes(self, blob: bytes, context: str, mode: Optional[str] = None) -> Tuple[Optional[bytes], Optional[str]]:
        mode = mode or self.image_mode
        try:
            image = Image.open(BytesIO(blob)).convert("RGB")
        except Exception:
            return None, None
        blocks = self._extract_image_blocks(image)
        if not blocks:
            return (blob if mode == "redesign" else None, None)

        notes = []
        redesigned = image.copy()
        draw = ImageDraw.Draw(redesigned)

        for block in blocks:
            original_text = block["text"]
            translated, _ = self._translate_text(original_text)
            if not translated:
                continue
            notes.append(f"{context}: {translated}")
            if mode == "redesign":
                self._draw_translated_block(draw, block["bbox"], translated)

        if mode == "redesign":
            buffer = BytesIO()
            redesigned.save(buffer, format="PNG")
            return buffer.getvalue(), None

        combined = "\n".join(notes) if notes else None
        return None, combined

    def _extract_image_blocks(self, image: Image.Image):
        try:
            data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT, lang=self.ocr_lang)
        except Exception as exc:
            LOGGER.warning("OCR error: %s", exc)
            return []
        blocks: Dict[int, Dict[str, any]] = {}
        for i, text in enumerate(data.get("text", [])):
            if not text or not text.strip():
                continue
            block_id = data["block_num"][i]
            left = data["left"][i]
            top = data["top"][i]
            width = data["width"][i]
            height = data["height"][i]
            entry = blocks.setdefault(
                block_id,
                {"bbox": [left, top, left + width, top + height], "text": []},
            )
            entry["bbox"][0] = min(entry["bbox"][0], left)
            entry["bbox"][1] = min(entry["bbox"][1], top)
            entry["bbox"][2] = max(entry["bbox"][2], left + width)
            entry["bbox"][3] = max(entry["bbox"][3], top + height)
            entry["text"].append(text.strip())
        results = []
        for entry in blocks.values():
            block_text = " ".join(entry["text"]).strip()
            if block_text:
                results.append({"bbox": entry["bbox"], "text": block_text})
        return results

    def _draw_translated_block(self, draw: ImageDraw.Draw, bbox: List[int], text: str) -> None:
        x0, y0, x1, y1 = bbox
        draw.rectangle([x0, y0, x1, y1], fill="white")
        height = max(12, int((y1 - y0) / max(1, len(text.splitlines()))))
        try:
            font = ImageFont.truetype("DejaVuSans.ttf", height)
        except Exception:
            font = ImageFont.load_default()
        draw.multiline_text((x0 + 2, y0 + 2), text, fill="black", font=font)

    def _process_pdf_images(self, page: fitz.Page, allow_replace: bool) -> None:
        images = page.get_images(full=True)
        if not images:
            return
        for idx, image in enumerate(images, start=1):
            xref = image[0]
            pix = fitz.Pixmap(page.parent, xref)
            if pix.n >= 5:
                pix = fitz.Pixmap(fitz.csRGB, pix)
            image_bytes = pix.tobytes("png")
            effective_mode = "notes" if (self.image_mode == "redesign" and not allow_replace) else self.image_mode
            new_blob, note = self._process_image_bytes(
                image_bytes, f"PDF страница {page.number + 1}, изображение {idx}", mode=effective_mode
            )
            if new_blob is not None and self.image_mode == "redesign" and allow_replace:
                page.parent.update_image(xref, stream=new_blob)
            if note:
                self.image_notes.append(note)

    def _process_docx_images(self, doc: Document) -> None:
        for shape in getattr(doc, "inline_shapes", []):
            try:
                image_rid = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
                image_part = doc.part.related_parts[image_rid]
                blob = image_part._blob
                new_blob, note = self._process_image_bytes(blob, "Документ")
                if new_blob is not None and self.image_mode == "redesign":
                    image_part._blob = new_blob
                if note:
                    self.image_notes.append(note)
            except Exception as exc:
                LOGGER.warning("Не удалось обработать изображение DOCX: %s", exc)

    def _append_image_notes_docx(self, doc: Document) -> None:
        if self.image_mode == "notes" and self.image_notes:
            doc.add_paragraph("")
            doc.add_paragraph("Перевод текста с изображений:", style=None)
            for note in self.image_notes:
                doc.add_paragraph(note)
            self.image_notes.clear()

    def translate_xlsx(self, input_path: str, output_path: str) -> None:
        """Перевод XLSX: текст в ячейках и комментарии к длинным переводам."""
        LOGGER.info("Перевод XLSX: %s", input_path)
        self.image_notes = []
        wb = load_workbook(input_path)
        if self.chunk_mode == "full":
            batch = TranslationBatch("xlsx")
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if isinstance(cell.value, str) and cell.value.strip():
                            batch.add(cell.value, lambda value, target_cell=cell: setattr(target_cell, "value", value))
            translations = self._translate_batch(batch, "xlsx")
            batch.apply(translations, self._fallback_translate_single)
        else:
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if isinstance(cell.value, str) and cell.value.strip():
                            translated, comment = self._translate_text(cell.value)
                            cell.value = translated
                            if comment:
                                cell.comment = Comment(text=comment, author="InsightAudio")
                # TODO: обработка изображений в XLSX при необходимости
        wb.save(output_path)

    def translate_pdf(self, input_path: str, output_path: str, reflow: bool = False) -> None:
        """Перевод PDF: reflow (замена текста в исходном макете) или плоский текстовый вывод."""
        LOGGER.info("Перевод PDF: %s", input_path)
        if reflow:
            doc = fitz.open(input_path)
            for page in doc:
                blocks = page.get_text("blocks")
                translated_blocks = []
                for block in blocks:
                    rect = fitz.Rect(block[0], block[1], block[2], block[3])
                    text = block[4].strip()
                    if not text:
                        continue
                    translated, _ = self._translate_text(text)
                    translated_blocks.append((rect, translated))
                for rect, translated in translated_blocks:
                    page.add_redact_annot(rect, fill=(1, 1, 1))
                if translated_blocks:
                    page.apply_redactions()
                for rect, translated in translated_blocks:
                    page.insert_textbox(rect, translated, fontsize=12, color=(0, 0, 0))
                self._process_pdf_images(page, allow_replace=True)
            doc.save(output_path)
        else:
            source = fitz.open(input_path)
            target = fitz.open()
            for idx, page in enumerate(source, start=1):
                text = page.get_text().strip()
                translated, _ = self._translate_text(text)
                new_page = target.new_page(width=612, height=792)
                new_page.insert_textbox(fitz.Rect(36, 36, 576, 756), translated or "", fontsize=11, color=(0, 0, 0))
                self._process_pdf_images(page, allow_replace=False)
            if self.image_notes and self.image_mode == "notes":
                notes_page = target.new_page(width=612, height=792)
                notes_page.insert_textbox(
                    fitz.Rect(36, 36, 576, 756),
                    "\n\n".join(self.image_notes),
                    fontsize=11,
                    color=(0, 0, 0),
                )
            target.save(output_path)
        self.image_notes.clear()

    def _translate_paragraph(self, paragraph):
        text = paragraph.text
        if not text.strip():
            return
        translated, comment = self._translate_text(text)
        paragraph.text = translated
        if comment:
            paragraph.add_run(f"\n{comment}")

    def _translate_text(self, text: str, instruction: str = "Переведи следующий текст") -> Tuple[str, Optional[str]]:
        """
        Центральная точка преобразования небольшого фрагмента текста.
        Здесь формируется prompt, делается вызов LLM и, при необходимости,
        возвращается не только переведённый текст, но и «короткая версия»
        (используется для заметок или ситуаций, когда перевод слишком длинный).
        """
        if not text.strip():
            return text, None
        instruction_text = (
            "Ты профессиональный переводчик. Переведи текст, сохранив смысл, структуру и числовые значения. "
            "Не добавляй пояснений, возвращай только перевод."
        )
        translated = self._call_translate_api(instruction_text, text)
        comment = None
        if len(translated) > len(text) * self.max_ratio:
            short_instruction = (
                f"Сократи перевод до объёма не больше чем {int(len(text) * self.max_ratio)} символов, "
                "сохранив ключевой смысл. Ответь только сокращённой версией."
            )
            short_version = self._call_translate_api(short_instruction, translated)
            comment = f"(Комментарий: полный перевод: \"{translated}\" | оригинал: \"{text}\")"
            translated = short_version
        return translated, comment

    def _call_translate_api(self, instruction: str, payload_text: str) -> str:
        """
        Оборачивает выбор backend’а: формирует итоговый prompt и делегирует вызов
        в Ollama / llama.cpp / пользовательский REST. Ожидает один логический блок.
        """
        if not payload_text.strip():
            return payload_text
        prompt = (
            f"{instruction}\n"
            f"Исходный язык: автоопределение\n"
            f"Целевой язык: {self.target_lang}\n"
            f"---\n"
            f"{payload_text.strip()}\n"
            f"---\n"
        )
        if self.backend == "ollama":
            return self._call_ollama(prompt)
        elif self.backend == "llama_cpp":
            return self._call_llama_cpp(prompt)
        elif self.backend == "custom_api":
            return self._call_custom_api(payload_text, prompt)
        return self._call_ollama(prompt)

    def _call_ollama(self, prompt: str) -> str:
        """
        Отправляет один блок текста в Ollama. Вызов идёт через streaming-утилиту,
        поэтому на выходе есть текст и замеренные метрики (ttft, tpot, throughput),
        которые потом попадают в UI и хранилище `model_metrics.json`.
        """
        payload = {
            "model": self.model,
            "prompt": prompt,
            "options": merge_model_options(self.model, {"temperature": 0.1}),
        }
        return generate_with_metrics(self.base_url, payload, self.model, scope="translate")

    def _call_llama_cpp(self, prompt: str) -> str:
        payload = {
            "model": self.model,
            "prompt": prompt,
            "n_predict": 512,
            "temperature": 0.1,
            "cache_prompt": True,
            "stream": False,
        }
        response = requests.post(self.base_url, json=payload, timeout=180)
        response.raise_for_status()
        data = response.json()
        if isinstance(data, dict):
            for key in ("content", "summary", "result", "output", "response"):
                if key in data and isinstance(data[key], str):
                    return data[key].strip()
        return json.dumps(data, ensure_ascii=False)

    def _call_custom_api(self, text: str, prompt: str) -> str:
        payload = {
            "model": self.model,
            "prompt": prompt,
            "text": text,
            "input": text,
            "target_language": self.target_lang,
        }
        response = requests.post(self.base_url, json=payload, timeout=180)
        response.raise_for_status()
        try:
            data = response.json()
        except ValueError:
            return response.text.strip()
        for key in ("translation", "result", "output", "text"):
            value = data.get(key)
            if isinstance(value, str):
                return value.strip()
        return json.dumps(data, ensure_ascii=False)


